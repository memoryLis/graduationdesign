#!/usr/bin/env python3
"""生成毕业论文答辩PPT —— 基于Spring Boot和Spring Cloud的商城系统的设计与实现"""

from pptx import Presentation
from pptx.util import Inches, Pt, Emu, Cm
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE, MSO_CONNECTOR_TYPE
from pptx.oxml.ns import qn
import os

# ── 全局配色（深蓝 + 白，学术风格）─────────────────────────────
DEEP_BLUE    = RGBColor(0x0A, 0x26, 0x47)   # 主色
MEDIUM_BLUE  = RGBColor(0x14, 0x42, 0x72)   # 辅助
LIGHT_BLUE   = RGBColor(0x2C, 0x6F, 0xAD)   # 强调
ACCENT_BLUE  = RGBColor(0x3A, 0x8E, 0xD8)   # 点缀
WHITE        = RGBColor(0xFF, 0xFF, 0xFF)
DARK_TEXT    = RGBColor(0x1A, 0x1A, 0x1A)
GRAY_TEXT    = RGBColor(0x80, 0x80, 0x80)
LIGHT_BG     = RGBColor(0xF0, 0xF3, 0xF8)   # 浅灰蓝背景
GOLD         = RGBColor(0xC9, 0xA0, 0x3E)   # 金色强调
TABLE_HDR    = RGBColor(0x0D, 0x33, 0x5A)   # 表头深色
TABLE_ROW1   = RGBColor(0xE8, 0xEE, 0xF4)   # 表格浅色行
RED_ACCENT   = RGBColor(0xC0, 0x39, 0x2B)   # 红色警示

prs = Presentation()
prs.slide_width  = Inches(13.333)
prs.slide_height = Inches(7.5)
W = prs.slide_width
H = prs.slide_height

# ── 工具函数 ─────────────────────────────────────────────────
def add_blank_slide():
    layout = prs.slide_layouts[6]  # blank
    return prs.slides.add_slide(layout)

def add_rect(slide, left, top, width, height, fill_color=None, border_color=None, border_width=None):
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, width, height)
    shape.line.fill.background()
    if fill_color:
        shape.fill.solid()
        shape.fill.fore_color.rgb = fill_color
    else:
        shape.fill.background()
    if border_color:
        shape.line.color.rgb = border_color
        shape.line.fill.solid()
        if border_width:
            shape.line.width = border_width
    return shape

def add_rounded_rect(slide, left, top, width, height, fill_color=None):
    shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height)
    shape.line.fill.background()
    if fill_color:
        shape.fill.solid()
        shape.fill.fore_color.rgb = fill_color
    else:
        shape.fill.background()
    return shape

def add_textbox(slide, left, top, width, height, text="", font_size=Pt(14),
                font_color=DARK_TEXT, bold=False, alignment=PP_ALIGN.LEFT,
                font_name="Microsoft YaHei", anchor=MSO_ANCHOR.TOP):
    txBox = slide.shapes.add_textbox(left, top, width, height)
    txBox.word_wrap = True
    tf = txBox.text_frame
    tf.word_wrap = True
    tf.auto_size = None
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = font_size
    p.font.color.rgb = font_color
    p.font.bold = bold
    p.font.name = font_name
    p.alignment = alignment
    # 东亚字体
    for run in p.runs:
        rPr = run._r.get_or_add_rPr()
        rPr.set(qn('a:eaTypeface'), font_name)
    return txBox

def add_multiline_textbox(slide, left, top, width, height, lines, font_name="Microsoft YaHei",
                          anchor=MSO_ANCHOR.TOP):
    """lines: list of (text, font_size, font_color, bold, alignment)"""
    txBox = slide.shapes.add_textbox(left, top, width, height)
    txBox.word_wrap = True
    tf = txBox.text_frame
    tf.word_wrap = True
    tf.auto_size = None
    for i, line_data in enumerate(lines):
        text, font_size, font_color, bold, alignment = line_data
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.text = text
        p.font.size = font_size
        p.font.color.rgb = font_color
        p.font.bold = bold
        p.font.name = font_name
        p.alignment = alignment
        p.space_after = Pt(4)
        for run in p.runs:
            rPr = run._r.get_or_add_rPr()
            rPr.set(qn('a:eaTypeface'), font_name)
    return txBox

def add_bottom_bar(slide):
    """底部装饰条"""
    add_rect(slide, Inches(0), H - Inches(0.08), W, Inches(0.08), fill_color=DEEP_BLUE)

def add_page_number(slide, num):
    """页码"""
    add_textbox(slide, W - Inches(1.2), H - Inches(0.55), Inches(1), Inches(0.4),
                str(num), Pt(10), GRAY_TEXT, alignment=PP_ALIGN.RIGHT)

def add_section_header(slide, title, subtitle=None):
    """统一的顶部标题区"""
    # 深蓝顶条
    add_rect(slide, Inches(0), Inches(0), W, Inches(1.25), fill_color=DEEP_BLUE)
    # 金色细线
    add_rect(slide, Inches(0.8), Inches(1.25), W - Inches(1.6), Pt(3), fill_color=GOLD)
    # 标题
    add_textbox(slide, Inches(0.8), Inches(0.15), W - Inches(1.6), Inches(0.7),
                title, Pt(30), WHITE, bold=True)
    if subtitle:
        add_textbox(slide, Inches(0.8), Inches(0.78), W - Inches(1.6), Inches(0.4),
                    subtitle, Pt(13), RGBColor(0xBB, 0xCC, 0xDD))
    add_bottom_bar(slide)

def add_icon_card(slide, left, top, width, height, icon_text, title, desc,
                  bg_color=WHITE, icon_bg=DEEP_BLUE):
    """带图标的卡片"""
    card = add_rounded_rect(slide, left, top, width, height, fill_color=bg_color)
    # 图标圆
    circle = slide.shapes.add_shape(MSO_SHAPE.OVAL, left + width/2 - Inches(0.45),
                                     top + Inches(0.25), Inches(0.9), Inches(0.9))
    circle.fill.solid()
    circle.fill.fore_color.rgb = icon_bg
    circle.line.fill.background()
    tf = circle.text_frame
    tf.word_wrap = False
    p = tf.paragraphs[0]
    p.text = icon_text
    p.font.size = Pt(22)
    p.font.color.rgb = WHITE
    p.alignment = PP_ALIGN.CENTER
    p.font.name = "Microsoft YaHei"
    # 标题
    add_textbox(slide, left + Inches(0.2), top + Inches(1.4), width - Inches(0.4), Inches(0.5),
                title, Pt(14), DEEP_BLUE, bold=True, alignment=PP_ALIGN.CENTER)
    # 描述
    add_textbox(slide, left + Inches(0.2), top + Inches(1.9), width - Inches(0.4), height - Inches(2.1),
                desc, Pt(10), GRAY_TEXT, alignment=PP_ALIGN.CENTER)

def add_table(slide, left, top, col_widths, headers, rows, header_font_size=Pt(12), cell_font_size=Pt(11)):
    """创建美观表格"""
    n_rows = len(rows) + 1
    n_cols = len(headers)
    total_w = sum(col_widths)
    table_shape = slide.shapes.add_table(n_rows, n_cols, left, top, total_w, Inches(0.45) * n_rows)
    table = table_shape.table

    for ci, cw in enumerate(col_widths):
        table.columns[ci].width = cw

    # 表头
    for ci, hdr in enumerate(headers):
        cell = table.cell(0, ci)
        cell.text = hdr
        cell.fill.solid()
        cell.fill.fore_color.rgb = TABLE_HDR
        for p in cell.text_frame.paragraphs:
            p.font.size = header_font_size
            p.font.color.rgb = WHITE
            p.font.bold = True
            p.font.name = "Microsoft YaHei"
            p.alignment = PP_ALIGN.CENTER

    # 数据行
    for ri, row in enumerate(rows):
        for ci, val in enumerate(row):
            cell = table.cell(ri + 1, ci)
            cell.text = str(val)
            cell.fill.solid()
            cell.fill.fore_color.rgb = TABLE_ROW1 if ri % 2 == 1 else WHITE
            for p in cell.text_frame.paragraphs:
                p.font.size = cell_font_size
                p.font.color.rgb = DARK_TEXT
                p.font.name = "Microsoft YaHei"
                p.alignment = PP_ALIGN.CENTER
    return table_shape

def add_flow_arrow(slide, left, top, width, height, text, fill=LIGHT_BLUE):
    """流程箭头框"""
    shape = slide.shapes.add_shape(MSO_SHAPE.CHEVRON, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill
    shape.line.fill.background()
    tf = shape.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(11)
    p.font.color.rgb = WHITE
    p.font.bold = True
    p.font.name = "Microsoft YaHei"
    p.alignment = PP_ALIGN.CENTER
    return shape

# ═══════════════════════════════════════════════════════════════
# 第1页：封面
# ═══════════════════════════════════════════════════════════════
slide1 = add_blank_slide()

# 深蓝背景
add_rect(slide1, Inches(0), Inches(0), W, H, fill_color=DEEP_BLUE)

# 装饰线条
add_rect(slide1, Inches(1.2), Inches(1.5), Inches(1.5), Pt(4), fill_color=GOLD)
add_rect(slide1, Inches(1.2), Inches(5.8), Inches(10.9), Pt(1.5), fill_color=GOLD)

# 校名
add_textbox(slide1, Inches(1.2), Inches(0.6), Inches(10), Inches(0.7),
            "井冈山大学", Pt(20), RGBColor(0xAA, 0xBB, 0xCC), bold=False)

# 副标题：毕业论文答辩
add_textbox(slide1, Inches(1.2), Inches(1.7), Inches(10), Inches(0.7),
            "本科毕业论文答辩", Pt(18), GOLD, bold=False)

# 论文题目
add_textbox(slide1, Inches(1.2), Inches(2.5), Inches(10.5), Inches(1.4),
            "基于Spring Boot和Spring Cloud的\n商城系统的设计与实现", Pt(36), WHITE, bold=True)

# 英文标题
add_textbox(slide1, Inches(1.2), Inches(4.1), Inches(10.5), Inches(0.6),
            "Design and Implementation of a Mall System Based on Spring Boot and Spring Cloud",
            Pt(12), RGBColor(0x99, 0xAA, 0xBB))

# 个人信息
info_lines = [
    ("答辩人：黄光亮", Pt(16), WHITE, False, PP_ALIGN.LEFT),
    ("学  号：2207104011", Pt(16), WHITE, False, PP_ALIGN.LEFT),
    ("专  业：数据科学与大数据技术", Pt(16), WHITE, False, PP_ALIGN.LEFT),
    ("学  院：电子与信息工程学院", Pt(16), WHITE, False, PP_ALIGN.LEFT),
    ("指导老师：曾宪文", Pt(16), WHITE, False, PP_ALIGN.LEFT),
]
add_multiline_textbox(slide1, Inches(1.2), Inches(5.2), Inches(6), Inches(2.0), info_lines)

# 学校logo占位区（右侧空白框）
logo_box = add_rounded_rect(slide1, Inches(9.8), Inches(1.0), Inches(2.8), Inches(2.8),
                             fill_color=RGBColor(0x0F, 0x30, 0x55))
add_textbox(slide1, Inches(9.8), Inches(2.0), Inches(2.8), Inches(0.8),
            "学校LOGO", Pt(14), RGBColor(0x88, 0x99, 0xAA), alignment=PP_ALIGN.CENTER)

# ═══════════════════════════════════════════════════════════════
# 第2页：目录
# ═══════════════════════════════════════════════════════════════
slide2 = add_blank_slide()
add_section_header(slide2, "目  录", "CONTENTS")
add_page_number(slide2, 2)

toc_items = [
    ("01", "选题原因", "研究背景与研究意义"),
    ("02", "研究内容", "系统核心功能模块与业务闭环"),
    ("03", "研究思路与方法", "技术路线与关键技术选型"),
    ("04", "研究成果", "系统架构实现与性能测试"),
    ("05", "总结与展望", "工作回顾与未来方向"),
]

for i, (num, title, desc) in enumerate(toc_items):
    y = Inches(1.8) + i * Inches(1.05)
    # 编号圆
    circle = slide2.shapes.add_shape(MSO_SHAPE.OVAL, Inches(1.2), y, Inches(0.65), Inches(0.65))
    circle.fill.solid()
    circle.fill.fore_color.rgb = DEEP_BLUE if i < 4 else LIGHT_BLUE
    circle.line.fill.background()
    tf = circle.text_frame
    p = tf.paragraphs[0]
    p.text = num
    p.font.size = Pt(16)
    p.font.color.rgb = WHITE
    p.font.bold = True
    p.font.name = "Microsoft YaHei"
    p.alignment = PP_ALIGN.CENTER
    # 连接线
    add_rect(slide2, Inches(2.05), y + Inches(0.3), Inches(0.5), Pt(2), fill_color=LIGHT_BLUE)
    # 标题
    add_textbox(slide2, Inches(2.7), y + Inches(0.02), Inches(8), Inches(0.4),
                title, Pt(22), DEEP_BLUE, bold=True)
    # 描述
    add_textbox(slide2, Inches(2.7), y + Inches(0.42), Inches(8), Inches(0.3),
                desc, Pt(12), GRAY_TEXT)

# ═══════════════════════════════════════════════════════════════
# 第3页：研究背景（选题原因-上）
# ═══════════════════════════════════════════════════════════════
slide3 = add_blank_slide()
add_section_header(slide3, "研究背景", "RESEARCH BACKGROUND —— 选题原因")
add_page_number(slide3, 3)

# 左侧：时代背景
bg_items = [
    ("🌐", "互联网电商蓬勃发展", "我国网上零售额持续增长，电子商务市场规模不断扩大，网上购物已成为人们日常消费方式"),
    ("⚠️", "传统单体架构困境", "代码耦合度高、扩展性差、维护成本高，高并发场景下系统崩溃，难以应对促销峰值流量"),
    ("🔧", "微服务架构兴起", "将单体应用拆分为小型独立服务，独立开发、部署、扩展，Spring Cloud成为主流微服务框架"),
    ("🇨🇳", "Spring Cloud Alibaba生态", "阿里巴巴基于双十一实战经验推出，提供Nacos、Sentinel等企业级组件"),
]

for i, (icon, title, desc) in enumerate(bg_items):
    y = Inches(1.7) + i * Inches(1.35)
    # 左侧色块标记
    add_rect(slide3, Inches(0.8), y, Pt(5), Inches(1.1), fill_color=DEEP_BLUE if i % 2 == 0 else LIGHT_BLUE)
    add_textbox(slide3, Inches(1.1), y, Inches(0.5), Inches(0.4), icon, Pt(22))
    add_textbox(slide3, Inches(1.7), y, Inches(5.5), Inches(0.35), title, Pt(16), DEEP_BLUE, bold=True)
    add_textbox(slide3, Inches(1.7), y + Inches(0.4), Inches(5.5), Inches(0.7), desc, Pt(11), GRAY_TEXT)

# 右侧：核心问题框
add_rounded_rect(slide3, Inches(8.0), Inches(1.7), Inches(4.8), Inches(5.0),
                 fill_color=LIGHT_BG)
add_textbox(slide3, Inches(8.3), Inches(1.9), Inches(4.2), Inches(0.4),
            "核心问题", Pt(18), DEEP_BLUE, bold=True, alignment=PP_ALIGN.CENTER)
add_rect(slide3, Inches(9.0), Inches(2.4), Inches(2.4), Pt(2), fill_color=GOLD)

problems = [
    "• 高并发访问导致系统崩溃",
    "• 业务耦合严重，维护困难",
    "• 部署效率低，扩展性差",
    "• 如何构建高可用的电商系统？",
]
add_multiline_textbox(slide3, Inches(8.3), Inches(2.7), Inches(4.2), Inches(3.5),
                      [(p, Pt(13), DARK_TEXT, False, PP_ALIGN.LEFT) for p in problems])

# ═══════════════════════════════════════════════════════════════
# 第4页：研究意义（选题原因-下）
# ═══════════════════════════════════════════════════════════════
slide4 = add_blank_slide()
add_section_header(slide4, "研究意义", "RESEARCH SIGNIFICANCE —— 选题原因")
add_page_number(slide4, 4)

significances = [
    ("📚", "理论意义",
     "结合Spring Cloud Alibaba技术生态，\n探索微服务架构在电商领域的应用实践，\n为同类系统开发提供参考方案",
     ["服务注册与发现", "分布式缓存策略", "异步消息解耦", "流量控制与熔断"]),
    ("🏭", "实践意义",
     "实现完整电商业务闭环，\n验证微服务架构在实际项目中的\n可行性和优越性",
     ["商品管理", "购物车与订单", "在线支付", "高并发验证"]),
    ("🎓", "学习意义",
     "涵盖当前主流Java开发技术栈，\n全面提升微服务架构\n设计与开发能力",
     ["Spring Boot", "Spring Cloud", "Vue3 + Vite", "Redis + MQ"]),
]

for i, (icon, title, desc, tags) in enumerate(significances):
    left = Inches(0.8) + i * Inches(4.1)
    # 卡片背景
    add_rounded_rect(slide4, left, Inches(1.8), Inches(3.8), Inches(5.0), fill_color=LIGHT_BG)
    # 图标
    add_textbox(slide4, left + Inches(0.3), Inches(2.0), Inches(3.2), Inches(0.5),
                icon + "  " + title, Pt(20), DEEP_BLUE, bold=True)
    add_rect(slide4, left + Inches(0.3), Inches(2.55), Inches(1.5), Pt(2), fill_color=GOLD)
    # 描述
    add_textbox(slide4, left + Inches(0.3), Inches(2.8), Inches(3.2), Inches(1.8),
                desc, Pt(12), DARK_TEXT)
    # 标签
    for j, tag in enumerate(tags):
        tag_left = left + Inches(0.3) + j * Inches(1.75) if j < 2 else left + Inches(0.3) + (j-2) * Inches(1.75)
        tag_top = Inches(4.7) if j < 2 else Inches(5.3)
        tag_shape = add_rounded_rect(slide4, tag_left, tag_top, Inches(1.55), Inches(0.42),
                                      fill_color=WHITE)
        tf = tag_shape.text_frame
        p = tf.paragraphs[0]
        p.text = tag
        p.font.size = Pt(9)
        p.font.color.rgb = MEDIUM_BLUE
        p.font.bold = True
        p.font.name = "Microsoft YaHei"
        p.alignment = PP_ALIGN.CENTER

# ═══════════════════════════════════════════════════════════════
# 第5页：研究内容
# ═══════════════════════════════════════════════════════════════
slide5 = add_blank_slide()
add_section_header(slide5, "研究内容", "RESEARCH CONTENT —— 论文框架与重点")
add_page_number(slide5, 5)

# 五大微服务模块
modules = [
    ("👤", "用户服务\nUser Service", "注册登录\nJWT鉴权\n地址管理", DEEP_BLUE),
    ("📦", "商品服务\nItem Service", "商品CRUD\n模糊搜索\n个性化推荐", MEDIUM_BLUE),
    ("🛒", "购物车服务\nCart Service", "购物车管理\n数量调整\n异步清理", LIGHT_BLUE),
    ("📋", "交易服务\nTrade Service", "订单创建\n状态管理\n超时取消", ACCENT_BLUE),
    ("💳", "支付服务\nPay Service", "余额支付\n支付通知\n状态同步", RGBColor(0x5A, 0xA8, 0xE0)),
]

for i, (icon, name, func, color) in enumerate(modules):
    left = Inches(0.5) + i * Inches(2.5)
    card = add_rounded_rect(slide5, left, Inches(1.8), Inches(2.3), Inches(3.2), fill_color=LIGHT_BG)
    # 顶部色条
    add_rect(slide5, left, Inches(1.8), Inches(2.3), Pt(5), fill_color=color)
    # 图标+名称
    add_textbox(slide5, left + Inches(0.2), Inches(2.0), Inches(1.9), Inches(0.9),
                icon, Pt(16), color, alignment=PP_ALIGN.CENTER)
    add_textbox(slide5, left + Inches(0.2), Inches(2.55), Inches(1.9), Inches(0.7),
                name, Pt(13), DARK_TEXT, bold=True, alignment=PP_ALIGN.CENTER)
    add_rect(slide5, left + Inches(0.5), Inches(3.2), Inches(1.3), Pt(1), fill_color=GOLD)
    add_textbox(slide5, left + Inches(0.2), Inches(3.35), Inches(1.9), Inches(1.3),
                func, Pt(11), GRAY_TEXT, alignment=PP_ALIGN.CENTER)

# 底部：完整业务闭环
add_textbox(slide5, Inches(0.8), Inches(5.3), Inches(12), Inches(0.4),
            "▎完整电商业务闭环", Pt(16), DEEP_BLUE, bold=True)

flow_steps = ["商品浏览", "加入购物车", "提交订单", "在线支付", "订单完成"]
for i, step in enumerate(flow_steps):
    left = Inches(1.5) + i * Inches(2.2)
    add_flow_arrow(slide5, left, Inches(5.85), Inches(1.9), Inches(0.55), step,
                   fill=[DEEP_BLUE, MEDIUM_BLUE, LIGHT_BLUE, ACCENT_BLUE, RGBColor(0x5A, 0xA8, 0xE0)][i])

# 下半部分补充信息
add_textbox(slide5, Inches(0.8), Inches(6.65), Inches(12), Inches(0.4),
            "技术支撑：Nacos服务注册 ｜ Gateway统一网关 ｜ OpenFeign远程调用 ｜ Redis缓存 ｜ RabbitMQ异步消息 ｜ MySQL持久化",
            Pt(10), GRAY_TEXT, alignment=PP_ALIGN.CENTER)

# ═══════════════════════════════════════════════════════════════
# 第6页：研究思路
# ═══════════════════════════════════════════════════════════════
slide6 = add_blank_slide()
add_section_header(slide6, "研究思路", "RESEARCH APPROACH —— 论文框架与重点")
add_page_number(slide6, 6)

# 四阶段流程
stages = [
    ("01", "系统分析", "可行性分析\n功能需求分析\n用例分析\n非功能需求", DEEP_BLUE),
    ("02", "总体设计", "系统架构设计\n功能结构设计\n业务流程设计\n数据库设计", MEDIUM_BLUE),
    ("03", "系统实现", "微服务开发\n前端页面构建\n中间件集成\n个性化推荐算法", LIGHT_BLUE),
    ("04", "系统测试", "功能测试\n性能测试\n并发压力测试\n结果分析", ACCENT_BLUE),
]

for i, (num, title, desc, color) in enumerate(stages):
    left = Inches(0.6) + i * Inches(3.15)
    # 阶段卡片
    add_rounded_rect(slide6, left, Inches(1.8), Inches(2.9), Inches(4.3), fill_color=LIGHT_BG)
    # 顶部编号条
    add_rect(slide6, left, Inches(1.8), Inches(2.9), Inches(0.7), fill_color=color)
    add_textbox(slide6, left + Inches(0.2), Inches(1.85), Inches(2.5), Inches(0.55),
                f" Phase {num} ", Pt(18), WHITE, bold=True)
    # 标题
    add_textbox(slide6, left + Inches(0.3), Inches(2.7), Inches(2.3), Inches(0.5),
                title, Pt(22), color, bold=True, alignment=PP_ALIGN.CENTER)
    add_rect(slide6, left + Inches(0.8), Inches(3.2), Inches(1.3), Pt(2), fill_color=GOLD)
    # 内容
    add_textbox(slide6, left + Inches(0.3), Inches(3.45), Inches(2.3), Inches(2.2),
                desc, Pt(13), DARK_TEXT, alignment=PP_ALIGN.CENTER)

# 箭头连接
for i in range(3):
    left = Inches(3.65) + i * Inches(3.15)
    arrow = slide6.shapes.add_shape(MSO_SHAPE.RIGHT_ARROW, left, Inches(3.6), Inches(0.35), Inches(0.35))
    arrow.fill.solid()
    arrow.fill.fore_color.rgb = GOLD
    arrow.line.fill.background()

# 底部：论文结构
add_textbox(slide6, Inches(0.8), Inches(6.45), Inches(12), Inches(0.4),
            "论文结构：第1章 绪论 → 第2章 相关技术概述 → 第3章 系统分析 → 第4章 总体设计 → 第5章 系统实现 → 第6章 系统测试 → 第7章 总结与展望",
            Pt(10), GRAY_TEXT, alignment=PP_ALIGN.CENTER)

# ═══════════════════════════════════════════════════════════════
# 第7页：研究方法与关键技术
# ═══════════════════════════════════════════════════════════════
slide7 = add_blank_slide()
add_section_header(slide7, "研究方法与关键技术", "RESEARCH METHODS & KEY TECHNOLOGIES —— 论文框架与重点")
add_page_number(slide7, 7)

# 左侧：研究方法
add_textbox(slide7, Inches(0.8), Inches(1.7), Inches(5.5), Inches(0.5),
            "▎研究方法", Pt(20), DEEP_BLUE, bold=True)

methods = [
    ("文献研究法", "查阅国内外微服务架构、电商系统相关文献，\n了解技术发展趋势和最佳实践"),
    ("系统设计法", "采用前后端分离B/S架构，微服务拆分，\n分库设计，模块化开发"),
    ("实验测试法", "使用JMeter进行并发压力测试，\n验证系统性能与稳定性"),
]
for i, (m_title, m_desc) in enumerate(methods):
    y = Inches(2.4) + i * Inches(1.3)
    add_rect(slide7, Inches(0.8), y, Pt(4), Inches(1.05), fill_color=DEEP_BLUE if i % 2 == 0 else LIGHT_BLUE)
    add_textbox(slide7, Inches(1.1), y, Inches(5.2), Inches(0.35), m_title, Pt(15), DEEP_BLUE, bold=True)
    add_textbox(slide7, Inches(1.1), y + Inches(0.4), Inches(5.2), Inches(0.65), m_desc, Pt(11), GRAY_TEXT)

# 右侧：核心技术栈
add_textbox(slide7, Inches(7.2), Inches(1.7), Inches(5.5), Inches(0.5),
            "▎核心技术栈", Pt(20), DEEP_BLUE, bold=True)

techs = [
    ("Spring Boot 2.7.12", "基础框架", DEEP_BLUE),
    ("Spring Cloud Alibaba", "微服务治理", MEDIUM_BLUE),
    ("Nacos", "注册 & 配置中心", LIGHT_BLUE),
    ("Spring Cloud Gateway", "API网关 & JWT鉴权", ACCENT_BLUE),
    ("OpenFeign + OkHttp", "服务间远程调用", RGBColor(0x5A, 0xA8, 0xE0)),
    ("Redis", "分布式缓存 & 分布式锁", RGBColor(0xDC, 0x56, 0x2C)),
    ("RabbitMQ", "异步消息 & 延迟队列", RGBColor(0xFF, 0x66, 0x00)),
    ("MySQL 8.0 + MyBatis-Plus", "数据持久化", RGBColor(0x44, 0x78, 0xA0)),
    ("Vue3 + Vite + Pinia", "前端框架", RGBColor(0x42, 0xB8, 0x83)),
]

for i, (name, role, color) in enumerate(techs):
    row = i // 3
    col = i % 3
    left = Inches(7.2) + col * Inches(2.0)
    y = Inches(2.4) + row * Inches(1.3)
    tag = add_rounded_rect(slide7, left, y, Inches(1.85), Inches(1.05), fill_color=LIGHT_BG)
    add_rect(slide7, left, y, Inches(1.85), Pt(4), fill_color=color)
    add_textbox(slide7, left + Inches(0.1), y + Inches(0.12), Inches(1.65), Inches(0.4),
                name, Pt(10), DARK_TEXT, bold=True, alignment=PP_ALIGN.CENTER)
    add_textbox(slide7, left + Inches(0.1), y + Inches(0.55), Inches(1.65), Inches(0.4),
                role, Pt(8), color, alignment=PP_ALIGN.CENTER)

# ═══════════════════════════════════════════════════════════════
# 第8页：研究成果 —— 系统架构
# ═══════════════════════════════════════════════════════════════
slide8 = add_blank_slide()
add_section_header(slide8, "研究成果 —— 系统架构", "RESEARCH RESULTS —— 论文框架与重点")
add_page_number(slide8, 8)

# 四层架构可视化
layers = [
    ("客户端层", "Vue3 Web应用\n浏览器访问\n响应式界面", DEEP_BLUE, Inches(1.8)),
    ("接入层", "Spring Cloud Gateway\n统一路由 · JWT鉴权\n动态路由配置", MEDIUM_BLUE, Inches(3.05)),
    ("服务层", "用户服务 | 商品服务 | 购物车服务\n交易服务 | 支付服务\nNacos · OpenFeign · Sentinel", LIGHT_BLUE, Inches(4.3)),
    ("数据层", "MySQL 数据库（分库）\nRedis 缓存\nRabbitMQ 消息队列", ACCENT_BLUE, Inches(5.55)),
]

for name, desc, color, y in layers:
    add_rounded_rect(slide8, Inches(0.8), y, Inches(6.5), Inches(1.05), fill_color=color)
    add_textbox(slide8, Inches(1.1), y + Inches(0.08), Inches(1.8), Inches(0.35),
                name, Pt(14), WHITE, bold=True)
    add_textbox(slide8, Inches(1.1), y + Inches(0.45), Inches(5.9), Inches(0.55),
                desc, Pt(10), RGBColor(0xDD, 0xE8, 0xF2))

# 箭头连接
for y in [Inches(2.9), Inches(4.15), Inches(5.4)]:
    arrow = slide8.shapes.add_shape(MSO_SHAPE.DOWN_ARROW, Inches(3.7), y, Inches(0.3), Inches(0.22))
    arrow.fill.solid()
    arrow.fill.fore_color.rgb = GOLD
    arrow.line.fill.background()

# 右侧：架构亮点
add_textbox(slide8, Inches(8.0), Inches(1.7), Inches(4.8), Inches(0.5),
            "▎架构亮点", Pt(18), DEEP_BLUE, bold=True)

highlights = [
    ("🎯", "前后端分离", "B/S架构，Vue3 + 网关 + 微服务"),
    ("🔀", "服务拆分", "按业务域拆分为5大独立微服务"),
    ("📡", "服务治理", "Nacos注册发现 + Gateway统一入口"),
    ("⚡", "异步解耦", "RabbitMQ处理支付通知/超时取消"),
    ("🚀", "高性能", "Redis缓存 + 乐观锁防超卖"),
    ("🔄", "个性化推荐", "基于余弦相似度的协同过滤算法"),
]
for i, (icon, title, desc) in enumerate(highlights):
    y = Inches(2.4) + i * Inches(0.78)
    add_textbox(slide8, Inches(8.0), y, Inches(0.4), Inches(0.35), icon, Pt(14))
    add_textbox(slide8, Inches(8.45), y, Inches(1.5), Inches(0.35), title, Pt(13), DEEP_BLUE, bold=True)
    add_textbox(slide8, Inches(10.0), y, Inches(3.0), Inches(0.35), desc, Pt(11), GRAY_TEXT)

# ═══════════════════════════════════════════════════════════════
# 第9页：研究成果 —— 性能测试
# ═══════════════════════════════════════════════════════════════
slide9 = add_blank_slide()
add_section_header(slide9, "研究成果 —— 性能测试", "RESEARCH RESULTS —— 论文框架与重点")
add_page_number(slide9, 9)

# 左侧：测试方案
add_textbox(slide9, Inches(0.8), Inches(1.7), Inches(5.5), Inches(0.5),
            "▎JMeter并发压力测试方案", Pt(18), DEEP_BLUE, bold=True)

add_textbox(slide9, Inches(0.8), Inches(2.3), Inches(5.8), Inches(1.0),
            "测试场景：模拟多用户同时对同一热门商品发起下单请求\n"
            "测试工具：Apache JMeter 5.x\n"
            "核心指标：平均响应时间、吞吐量、错误率、库存一致性",
            Pt(12), DARK_TEXT)

# 测试方案表
test_plan_headers = ["测试组", "并发用户数", "循环次数", "总请求数"]
test_plan_rows = [
    ["第一组", "10", "10", "100"],
    ["第二组", "20", "10", "200"],
    ["第三组", "50", "10", "500"],
]
add_table(slide9, Inches(0.8), Inches(3.5),
          [Inches(1.4), Inches(1.5), Inches(1.5), Inches(1.5)],
          test_plan_headers, test_plan_rows)

# 右侧：测试结果
add_textbox(slide9, Inches(7.2), Inches(1.7), Inches(5.5), Inches(0.5),
            "▎性能测试结果", Pt(18), DEEP_BLUE, bold=True)

result_headers = ["并发数", "平均响应(ms)", "吞吐量(req/s)", "错误率", "库存一致性"]
result_rows = [
    ["100", "156", "320", "0%", "一致 ✓"],
    ["200", "287", "285", "0%", "一致 ✓"],
    ["500", "512", "245", "0.2%", "一致 ✓"],
]
add_table(slide9, Inches(7.2), Inches(2.3),
          [Inches(1.2), Inches(1.4), Inches(1.4), Inches(1.1), Inches(1.3)],
          result_headers, result_rows)

# 结论卡片
conclusions = [
    ("✓ 100并发", "响应156ms\n零错误率", DEEP_BLUE),
    ("✓ 500并发", "响应512ms\n错误率仅0.2%", MEDIUM_BLUE),
    ("✓ 库存安全", "所有测试组\n零超卖", LIGHT_BLUE),
]
for i, (title, desc, color) in enumerate(conclusions):
    left = Inches(7.2) + i * Inches(2.1)
    card = add_rounded_rect(slide9, left, Inches(4.9), Inches(1.9), Inches(1.6), fill_color=color)
    add_textbox(slide9, left + Inches(0.1), Inches(5.05), Inches(1.7), Inches(0.4),
                title, Pt(14), WHITE, bold=True, alignment=PP_ALIGN.CENTER)
    add_textbox(slide9, left + Inches(0.1), Inches(5.45), Inches(1.7), Inches(0.8),
                desc, Pt(11), RGBColor(0xDD, 0xE8, 0xF2), alignment=PP_ALIGN.CENTER)

# ═══════════════════════════════════════════════════════════════
# 第10页：总结与展望
# ═══════════════════════════════════════════════════════════════
slide10 = add_blank_slide()
add_section_header(slide10, "总结与展望", "SUMMARY & OUTLOOK")
add_page_number(slide10, 10)

# 左侧：工作总结
add_textbox(slide10, Inches(0.8), Inches(1.7), Inches(6.0), Inches(0.5),
            "▎工作总结", Pt(22), DEEP_BLUE, bold=True)

summary_items = [
    ("01", "完成微服务架构设计",
     "基于Spring Cloud Alibaba，实现五大核心微服务模块，通过Nacos、Gateway、OpenFeign构建完整微服务治理体系"),
    ("02", "实现完整电商业务闭环",
     "涵盖商品浏览搜索、购物车管理、订单交易、在线支付、用户管理、后台管理等核心功能"),
    ("03", "多技术手段保障性能",
     "Redis缓存提升查询性能、RabbitMQ异步解耦、乐观锁防超卖、个性化推荐算法"),
    ("04", "JMeter验证高并发能力",
     "500并发下响应512ms，错误率仅0.2%，零超卖，系统稳定可靠"),
]

for i, (num, title, desc) in enumerate(summary_items):
    y = Inches(2.4) + i * Inches(1.15)
    add_rect(slide10, Inches(0.8), y, Pt(4), Inches(0.9), fill_color=DEEP_BLUE if i % 2 == 0 else LIGHT_BLUE)
    add_textbox(slide10, Inches(1.1), y, Inches(0.5), Inches(0.3), num, Pt(18), GOLD, bold=True)
    add_textbox(slide10, Inches(1.6), y, Inches(5.2), Inches(0.3), title, Pt(15), DEEP_BLUE, bold=True)
    add_textbox(slide10, Inches(1.6), y + Inches(0.38), Inches(5.2), Inches(0.65), desc, Pt(10), GRAY_TEXT)

# 右侧：未来展望
add_textbox(slide10, Inches(7.2), Inches(1.7), Inches(5.5), Inches(0.5),
            "▎未来展望", Pt(22), DEEP_BLUE, bold=True)

outlooks = [
    ("🐳", "容器化部署", "引入Docker + Kubernetes，\n实现自动化部署与弹性伸缩"),
    ("🔍", "全链路追踪", "引入SkyWalking/Zipkin，\n实现调用链路可视化监控"),
    ("📊", "读写分离", "MySQL主从复制 + 读写分离，\n提升数据层性能与可用性"),
    ("🌐", "Service Mesh", "引入Istio服务网格，\n下沉服务治理能力至基础设施层"),
]
for i, (icon, title, desc) in enumerate(outlooks):
    y = Inches(2.4) + i * Inches(1.2)
    card = add_rounded_rect(slide10, Inches(7.2), y, Inches(5.3), Inches(1.0), fill_color=LIGHT_BG)
    add_rect(slide10, Inches(7.2), y, Pt(4), Inches(1.0), fill_color=ACCENT_BLUE)
    add_textbox(slide10, Inches(7.5), y + Inches(0.08), Inches(0.5), Inches(0.35), icon, Pt(16))
    add_textbox(slide10, Inches(8.1), y + Inches(0.08), Inches(1.8), Inches(0.35), title, Pt(14), DEEP_BLUE, bold=True)
    add_textbox(slide10, Inches(8.1), y + Inches(0.45), Inches(4.2), Inches(0.5), desc, Pt(10), GRAY_TEXT)

# ═══════════════════════════════════════════════════════════════
# 第11页：结束页
# ═══════════════════════════════════════════════════════════════
slide11 = add_blank_slide()

# 深蓝全背景
add_rect(slide11, Inches(0), Inches(0), W, H, fill_color=DEEP_BLUE)

# 装饰线条
add_rect(slide11, Inches(1.2), Inches(2.0), Inches(10.9), Pt(2), fill_color=GOLD)
add_rect(slide11, Inches(1.2), Inches(5.5), Inches(10.9), Pt(2), fill_color=GOLD)

# 感谢文字
add_textbox(slide11, Inches(1.2), Inches(2.5), Inches(10.9), Inches(1.2),
            "敬请老师批评指正", Pt(48), WHITE, bold=True, alignment=PP_ALIGN.CENTER)

add_textbox(slide11, Inches(1.2), Inches(3.8), Inches(10.9), Inches(0.7),
            "Thank you for your valuable comments and suggestions",
            Pt(18), RGBColor(0x99, 0xAA, 0xBB), alignment=PP_ALIGN.CENTER)

# 底部信息
bottom_lines = [
    ("答辩人：黄光亮", Pt(16), RGBColor(0xCC, 0xDD, 0xEE), False, PP_ALIGN.CENTER),
    ("指导老师：曾宪文", Pt(16), RGBColor(0xCC, 0xDD, 0xEE), False, PP_ALIGN.CENTER),
    ("井冈山大学 · 电子与信息工程学院", Pt(14), GOLD, False, PP_ALIGN.CENTER),
]
add_multiline_textbox(slide11, Inches(1.2), Inches(4.5), Inches(10.9), Inches(1.2), bottom_lines)

add_page_number(slide11, 11)

# ── 保存 ─────────────────────────────────────────────────────
output_path = "F:/graduationdesign/毕业论文答辩PPT.pptx"
prs.save(output_path)
print(f"PPT 已生成：{output_path}")
print(f"共 {len(prs.slides)} 页")
